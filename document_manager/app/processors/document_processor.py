import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import hashlib
from datetime import datetime
from typing import Dict, Any, List
import json

class DocumentProcessor:
    def __init__(self, config: dict):
        self.config = config
        
    def process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata and content from PDF files"""
        metadata = {
            'type': 'pdf',
            'processed_at': datetime.now().isoformat()
        }
        
        try:
            with pdfplumber.open(file_path) as pdf:
                metadata['pages'] = len(pdf.pages)
                metadata['metadata'] = pdf.metadata
                
                # Extract text from first 3 pages for classification
                text_content = []
                for i, page in enumerate(pdf.pages[:3]):
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
                
                metadata['sample_text'] = ' '.join(text_content)[:2000]
                
                # Extract tables if present
                tables = []
                for page in pdf.pages[:5]:  # Check first 5 pages
                    page_tables = page.extract_tables()
                    if page_tables:
                        tables.extend(page_tables)
                
                metadata['has_tables'] = len(tables) > 0
                metadata['table_count'] = len(tables)
                
                # Calculate file hash for deduplication
                with open(file_path, 'rb') as f:
                    metadata['hash'] = hashlib.md5(f.read()).hexdigest()
                    
        except Exception as e:
            metadata['error'] = str(e)
            
        return metadata
    
    def process_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from Word documents"""
        metadata = {
            'type': 'docx',
            'processed_at': datetime.now().isoformat()
        }
        
        try:
            doc = docx.Document(file_path)
            
            # Document properties
            metadata['properties'] = {
                'author': doc.core_properties.author,
                'created': str(doc.core_properties.created),
                'modified': str(doc.core_properties.modified),
                'title': doc.core_properties.title,
                'subject': doc.core_properties.subject
            }
            
            # Extract text
            text_content = []
            for paragraph in doc.paragraphs[:50]:  # First 50 paragraphs
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            metadata['sample_text'] = ' '.join(text_content)[:2000]
            metadata['paragraph_count'] = len(doc.paragraphs)
            metadata['table_count'] = len(doc.tables)
            
            # Extract headers
            headers = []
            for paragraph in doc.paragraphs:
                if paragraph.style.name.startswith('Heading'):
                    headers.append({
                        'level': paragraph.style.name,
                        'text': paragraph.text[:100]
                    })
            metadata['headers'] = headers[:10]  # First 10 headers
            
        except Exception as e:
            metadata['error'] = str(e)
            
        return metadata
    
    def process_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from Excel files"""
        metadata = {
            'type': 'xlsx',
            'processed_at': datetime.now().isoformat()
        }
        
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            
            metadata['sheet_names'] = wb.sheetnames
            metadata['sheet_count'] = len(wb.sheetnames)
            
            # Analyze first sheet
            if wb.sheetnames:
                sheet = wb[wb.sheetnames[0]]
                metadata['first_sheet_dimensions'] = {
                    'rows': sheet.max_row,
                    'columns': sheet.max_column
                }
                
                # Sample data from first sheet
                sample_data = []
                for row in sheet.iter_rows(max_row=5, values_only=True):
                    sample_data.append([str(cell)[:50] if cell else None 
                                       for cell in row[:5]])
                metadata['sample_data'] = sample_data
                
        except Exception as e:
            metadata['error'] = str(e)
            
        return metadata
    
    def process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from PowerPoint files"""
        metadata = {
            'type': 'pptx',
            'processed_at': datetime.now().isoformat()
        }
        
        try:
            prs = Presentation(file_path)
            
            metadata['slide_count'] = len(prs.slides)
            
            # Extract text from slides
            text_content = []
            slide_titles = []
            
            for slide in prs.slides[:10]:  # First 10 slides
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                            
                            # Check if it's a title
                            if shape == slide.shapes.title:
                                slide_titles.append(text)
                
                if slide_text:
                    text_content.append(' '.join(slide_text))
            
            metadata['sample_text'] = ' '.join(text_content)[:2000]
            metadata['slide_titles'] = slide_titles
            
        except Exception as e:
            metadata['error'] = str(e)
            
        return metadata
    
    def process_document(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Main method to process any document type"""
        processors = {
            'pdf': self.process_pdf,
            'docx': self.process_docx,
            'doc': self.process_docx,
            'xlsx': self.process_xlsx,
            'xls': self.process_xlsx,
            'pptx': self.process_pptx,
            'ppt': self.process_pptx
        }
        
        processor = processors.get(file_type.lower())
        if processor:
            return processor(file_path)
        else:
            return {'type': file_type, 'error': 'Unsupported file type'}
